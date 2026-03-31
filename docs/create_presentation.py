"""Generate I/O Optimization presentation from the plan (6 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Colors -------------------------------------------------------------------

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_ORANGE = RGBColor(0xFF, 0xA7, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
TABLE_HEADER_BG = RGBColor(0x2A, 0x2A, 0x45)
TABLE_ROW_BG = RGBColor(0x22, 0x22, 0x3A)
TABLE_ALT_BG = RGBColor(0x28, 0x28, 0x42)

# -- Helpers ------------------------------------------------------------------

L_MARGIN = Inches(0.6)
CONTENT_W = Inches(8.8)

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_paragraph(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(4), space_after=Pt(1), level=0):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_bullet(tf, text, size=12, color=WHITE, bold=False, level=0):
    return add_paragraph(tf, text, size=size, color=color, bold=bold, level=level)

def make_title(slide, title_text, subtitle_text=None):
    set_slide_bg(slide, DARK_BG)
    tb = add_textbox(slide, L_MARGIN, Inches(0.25), CONTENT_W, Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, title_text, size=24, color=ACCENT_BLUE, bold=True)
    if subtitle_text:
        add_paragraph(tf, subtitle_text, size=12, color=LIGHT_GRAY, space_before=Pt(6))

def add_table(slide, rows, cols, left, top, width, height, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            if r_idx == 0:
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = ACCENT_BLUE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                run.font.size = Pt(9)
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG if r_idx % 2 == 1 else TABLE_ALT_BG
            p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
    return table_shape

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# -- Build presentation -------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

# =============================================================================
# SLIDE 1 — Why Offload KV Cache to Storage?
# =============================================================================
s1 = prs.slides.add_slide(BLANK)
make_title(s1, "Why Offload KV Cache to Storage?", "Optimizing the CPU \u2194 Storage leg for LLMD")

# Three-tier diagram
tier_top = Inches(1.05)
tier_h = Inches(0.7)
tier_w = Inches(2.6)
tier_gap = Inches(0.3)
tier_labels = [("GPU memory", "fastest, smallest"), ("CPU memory", "fast, larger"), ("Storage", "largest, persistent")]
tier_colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE]

for i, ((label, sub), clr) in enumerate(zip(tier_labels, tier_colors)):
    left = L_MARGIN + i * (tier_w + tier_gap)
    shape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, tier_top, tier_w, tier_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)
    tf_s = shape.text_frame
    tf_s.word_wrap = True
    tf_s.margin_top = Pt(4)
    tf_s.margin_bottom = Pt(2)
    p1 = tf_s.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(11)
    r1.font.color.rgb = clr
    r1.font.bold = True
    p2 = tf_s.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(8)
    r2.font.color.rgb = LIGHT_GRAY
    # Arrow
    if i < 2:
        aleft = left + tier_w
        atop = tier_top + tier_h / 2 - Pt(7)
        atb = add_textbox(s1, aleft, atop, tier_gap, Inches(0.25))
        ap = atb.text_frame.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = "\u2192"
        ar.font.size = Pt(16)
        ar.font.color.rgb = LIGHT_GRAY

# Scope labels under diagram
tb_scope = add_textbox(s1, L_MARGIN, Inches(1.8), CONTENT_W, Inches(0.35))
tf_scope = tb_scope.text_frame
tf_scope.word_wrap = True
set_text(tf_scope, "GPU \u2192 CPU: already in LLMD          CPU \u2194 Storage: what we\u2019re optimizing", size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Left column: Why not GPU → Storage directly?
tb_left = add_textbox(s1, L_MARGIN, Inches(2.2), Inches(4.6), Inches(3.0))
tf_left = tb_left.text_frame
tf_left.word_wrap = True
set_text(tf_left, "GPU \u2192 Storage directly (GDS) \u2014 preferred, but often unavailable:", size=12, color=ACCENT_ORANGE, bold=True)
add_bullet(tf_left, "GDS is the fastest path \u2014 bypasses CPU entirely via DMA/RDMA", size=10, color=ACCENT_GREEN)
add_bullet(tf_left, "But: requires specific NVIDIA hardware + NIC + storage drivers", size=10, color=WHITE)
add_bullet(tf_left, "Works best with local NVMe; on shared FS often unsupported or falls back to CPU", size=10, color=WHITE)
add_bullet(tf_left, "Ties up GPU during transfer \u2014 can\u2019t do inference while doing I/O", size=10, color=WHITE)
add_bullet(tf_left, "Not available in all environments (cloud, heterogeneous clusters)", size=10, color=WHITE)
add_paragraph(tf_left, "When GDS isn\u2019t available \u2192 the CPU path is the fallback:", size=11, color=ACCENT_BLUE, bold=True, space_before=Pt(8))
add_bullet(tf_left, "Works with any storage, any hardware, any deployment", size=10, color=WHITE)
add_bullet(tf_left, "GPU stays free for inference while CPU handles I/O", size=10, color=WHITE)

# Right column: scenarios
tb_right = add_textbox(s1, Inches(5.2), Inches(2.2), Inches(4.2), Inches(3.0))
tf_right = tb_right.text_frame
tf_right.word_wrap = True
set_text(tf_right, "Key scenarios:", size=13, color=ACCENT_GREEN, bold=True)
add_bullet(tf_right, "1. New decode instance \u2014 load from storage, skip recalculation", size=11, color=WHITE)
add_bullet(tf_right, "2. Cache evicted from GPUs \u2014 persist to storage, reload later", size=11, color=WHITE)
add_bullet(tf_right, "3. Preemption / migration \u2014 persist cache, resume on another GPU", size=11, color=WHITE)
add_paragraph(tf_right, "Read path (storage \u2192 CPU) is latency-critical", size=10, color=ACCENT_ORANGE, bold=True, space_before=Pt(10))
add_paragraph(tf_right, "Write path can be backgrounded during inference", size=10, color=LIGHT_GRAY)

add_notes(s1, """SPEAKER NOTES - Slide 1: Why Offload KV Cache to Storage?

- The core tradeoff: recalculating KV cache requires a full forward pass through the model (prefill), which is expensive in both compute and latency. If we can reload from storage faster than recalculating, we save GPU cycles and reduce time-to-first-token.

- The three-tier hierarchy: GPU -> CPU eviction already exists in LLMD and is fast (PCIe/NVLink bandwidth). Our project focuses on the second leg: CPU <-> Storage. This leg is triggered when: (a) CPU memory itself is under pressure and must overflow, or (b) the cache needs to persist beyond process lifetime (restarts, migrations, failures).

- CPU memory as hot cache: In many cases, KV blocks evicted from GPU live in CPU memory and are reloaded to GPU without ever touching storage. Storage is the cold tier -- only needed when CPU memory is full or when persistence is required. This means the CPU->Storage write path can often be backgrounded (not urgent), but the Storage->CPU read path is latency-critical when it does happen (cache miss in CPU memory).

- GPU -> Storage directly (GDS) -- preferred when available, but often isn't:
  GPUDirect Storage (GDS) is the fastest path -- bypasses CPU entirely via DMA/RDMA. When available, it's preferred. But it has significant availability limitations:
  - Hardware requirements: GDS requires a specific NVIDIA GPU + Mellanox NIC + compatible storage controller + MLNX_OFED drivers. Many production deployments don't have the full stack.
  - Filesystem compatibility: GDS works best with local NVMe. On network/shared filesystems (Spectrum Scale, Lustre, GPFS), GDS is often unsupported or silently falls back to the CPU-bounce path anyway -- giving you worse performance than a well-optimized CPU path because of the fallback overhead.
  - GPU utilization tradeoff: While GDS transfers data, the GPU's DMA engine is busy. For inference serving, every microsecond of GPU time matters. The CPU path lets the GPU keep doing inference while CPU and storage handle I/O in the background -- better end-to-end throughput.
  - Portability: Different clouds, on-prem clusters, and edge deployments have different hardware. The CPU path works everywhere with zero special requirements -- just Linux + any filesystem.

- In practice, both paths will coexist. GDS for specific cases where hardware is available and latency is critical (e.g., local NVMe on a DGX). CPU path as the general-purpose layer that works across all deployments. Optimizing the CPU path has broad impact because it applies everywhere, not just on specific hardware configurations.

- Additional scenarios (mention if asked): long-context conversations (KV cache grows beyond GPU memory, offload older tokens), batch size scaling (many concurrent requests overflow to storage).

- In all scenarios, the read path (storage -> CPU -> GPU) is on the critical path -- it directly impacts latency the user sees. Write path can be async/backgrounded.""")


# =============================================================================
# SLIDE 2 — Benchmark Project
# =============================================================================
s2 = prs.slides.add_slide(BLANK)
make_title(s2, "Benchmark Project")

tb2 = add_textbox(s2, L_MARGIN, Inches(0.95), CONTENT_W, Inches(4.3))
tf2 = tb2.text_frame
tf2.word_wrap = True
set_text(tf2, "Framework by [team name] for evaluating I/O strategies", size=13, color=LIGHT_GRAY)

add_paragraph(tf2, "4 backends tested:", size=13, color=ACCENT_GREEN, bold=True, space_before=Pt(10))
add_bullet(tf2, "C++ \u2014 pread/pwrite + custom thread pool, GIL released, serialized file creation", size=11, color=WHITE)
add_bullet(tf2, "Python self \u2014 asyncio + ThreadPoolExecutor + os.readv/write", size=11, color=WHITE)
add_bullet(tf2, "aiofiles \u2014 pure Python async I/O library (cross-platform baseline)", size=11, color=WHITE)
add_bullet(tf2, "NIXL \u2014 NVIDIA library with POSIX backend, memory registration, descriptor-based transfers", size=11, color=WHITE)

add_paragraph(tf2, "Test matrix:", size=13, color=ACCENT_GREEN, bold=True, space_before=Pt(10))
add_bullet(tf2, "Block sizes: 2, 4, 8, 16, 32, 64 MB (KV cache blocks are ~5-50MB depending on model)", size=11, color=WHITE)
add_bullet(tf2, "Thread counts: 16, 32, 64", size=11, color=WHITE)
add_bullet(tf2, "Modes: sequential write/read, fixed total data (100GB), concurrent read+write", size=11, color=WHITE)

add_paragraph(tf2, "Environment:", size=13, color=ACCENT_GREEN, bold=True, space_before=Pt(10))
add_bullet(tf2, "K8s pods, 450Gi RAM, /dev/shm (tmpfs) + persistent storage", size=11, color=WHITE)
add_bullet(tf2, "All backends use atomic writes (temp file + rename) and cache cleaning between phases", size=11, color=WHITE)

add_notes(s2, """SPEAKER NOTES - Slide 2: Benchmark Project

- Credit the other team's work -- solid foundation, well-structured benchmarks.

- Block size depends on model: LLaMA-70B ~52MB, LLaMA-13B ~13MB, quantized ~3-7MB. The tested range (2-64MB) covers realistic KV cache block sizes.

- Cache cleaning between phases: reads 100GB of unrelated data to evict written data from OS page cache. This ensures we measure true storage throughput, not cached reads. Critical for meaningful results.

- Atomic write pattern (temp file + rename): prevents partial/corrupt files if process crashes mid-write. All 4 backends implement this.

- The benchmark framework is reusable -- we'll add new backends and parameters to it for our optimization work.

- NIXL (NVIDIA Interconnect eXchange Library) -- an NVIDIA library designed for high-performance data transfers, primarily GPU-to-GPU and RDMA across nodes. It supports a POSIX file backend, which is what's tested here. The library uses explicit memory registration (pinning regions for DMA) and a descriptor-based transfer model. File I/O isn't its primary use case -- it's optimized for network/GPU transfers. The busy-polling completion model also wastes CPU cycles compared to blocking reads. We included it because it's being evaluated within the LLMD ecosystem for other transfer paths.

- aiofiles -- a pure Python async I/O library. Cross-platform, easy to use, but adds abstraction overhead on top of the OS. Included as a baseline to show the cost of high-level abstractions vs native implementations. Not expected to compete with C++ or even the Python self-implementation.""")


# =============================================================================
# SLIDE 3 — Current Results (baseline)
# =============================================================================
s3 = prs.slides.add_slide(BLANK)
make_title(s3, "Current Results \u2014 Baseline")

data3 = [
    ["Backend", "Peak Read", "Peak Write", "Notes"],
    ["C++ (pread + threadpool)", "~55-60 GB/s", "~40-45 GB/s", "Best overall"],
    ["Python self (readv)", "~45-55 GB/s", "~25-35 GB/s", "GIL limits writes"],
    ["NIXL", "~15-20 GB/s", "~10-15 GB/s", "Busy-poll overhead"],
]
add_table(s3, 4, 4, L_MARGIN, Inches(1.0), CONTENT_W, Inches(1.15), data3,
          col_widths=[Inches(2.4), Inches(1.8), Inches(1.8), Inches(2.8)])

tb3 = add_textbox(s3, L_MARGIN, Inches(2.35), CONTENT_W, Inches(2.8))
tf3 = tb3.text_frame
tf3.word_wrap = True
set_text(tf3, "Key Observations", size=14, color=ACCENT_GREEN, bold=True)
add_bullet(tf3, "Throughput plateaus at 16-32MB block size", size=11, color=WHITE)
add_bullet(tf3, "Diminishing returns beyond 32 threads", size=11, color=WHITE)
add_bullet(tf3, "C++ wins: GIL released, native thread pool, serialized file creation avoids dir contention", size=11, color=WHITE)
add_bullet(tf3, "Tested on /dev/shm (tmpfs) \u2014 100GB data, 5 iterations, cache cleaned between phases", size=10, color=LIGHT_GRAY)

add_paragraph(tf3, "This is our baseline \u2014 the C++ backend using standard POSIX pread/pwrite",
              size=12, color=ACCENT_ORANGE, bold=True, space_before=Pt(10))

add_notes(s3, """SPEAKER NOTES - Slide 3: Current Results (baseline)

- tmpfs (/dev/shm) -- a filesystem that lives entirely in RAM, not on any physical storage device. Reads and writes go directly to/from memory, no disk involved. This gives us the theoretical upper bound for I/O throughput -- any real storage (NVMe, Spectrum Scale) will be slower. We benchmark on tmpfs first to measure software overhead in isolation, without storage hardware being the bottleneck. On real storage, the same software optimizations will have proportionally larger impact since there's more overhead to eliminate.

- Explain why C++ beats Python: GIL release during I/O, native thread pool vs Python's ThreadPoolExecutor overhead.

- NIXL underperforms here because file I/O isn't its primary use case -- it's optimized for GPU-to-GPU/RDMA transfers. See slide 2 notes for details.

- The plateau at 16-32MB suggests we're hitting syscall overhead and page cache copy limits, not device bandwidth -- room to improve.

- These are averages over 5 iterations with 100GB cache cleaning between write and read phases.""")


# =============================================================================
# SLIDE 4 — How Linux I/O Works
# =============================================================================
s4 = prs.slides.add_slide(BLANK)
make_title(s4, "How Linux I/O Works \u2014 Where We Can Improve")

# Data path diagram
diag_top = Inches(1.05)
box_h = Inches(0.48)
box_w = Inches(1.3)
gap = Inches(0.18)

labels = ["User Buffer", "syscall\n(ctx switch)", "Page Cache\n(kernel copy)", "Block Layer\n(scheduler)", "Device Driver", "Storage"]
diag_colors = [ACCENT_BLUE, ACCENT_ORANGE, ACCENT_ORANGE, ACCENT_ORANGE, LIGHT_GRAY, LIGHT_GRAY]

for i, (label, clr) in enumerate(zip(labels, diag_colors)):
    left = L_MARGIN + i * (box_w + gap)
    shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, diag_top, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)
    tf_s = shape.text_frame
    tf_s.word_wrap = True
    tf_s.margin_top = Pt(2)
    tf_s.margin_bottom = Pt(2)
    p = tf_s.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(8)
    run.font.color.rgb = clr
    run.font.bold = True
    if i < len(labels) - 1:
        aleft = left + box_w
        atop = diag_top + box_h / 2 - Pt(6)
        atb = add_textbox(s4, aleft, atop, gap, Inches(0.2))
        ap = atb.text_frame.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = "\u2192"
        ar.font.size = Pt(11)
        ar.font.color.rgb = LIGHT_GRAY

# Bottlenecks (left)
tb4l = add_textbox(s4, L_MARGIN, Inches(1.8), Inches(4.0), Inches(2.2))
tf4l = tb4l.text_frame
tf4l.word_wrap = True
set_text(tf4l, "Three bottlenecks:", size=13, color=ACCENT_BLUE, bold=True)
add_bullet(tf4l, "1. Syscall overhead \u2014 1 context switch per pread, \u00d71000+ blocks", size=11, color=WHITE)
add_bullet(tf4l, "2. Page cache copy \u2014 double copy through kernel buffers", size=11, color=WHITE)
add_bullet(tf4l, "3. Scheduling \u2014 kernel doesn\u2019t know our access pattern", size=11, color=WHITE)

# Solutions (right)
tb4r = add_textbox(s4, Inches(5.0), Inches(1.8), Inches(4.4), Inches(3.4))
tf4r = tb4r.text_frame
tf4r.word_wrap = True
set_text(tf4r, "What eliminates each:", size=13, color=ACCENT_GREEN, bold=True)

add_bullet(tf4r, "io_uring \u2192 eliminates #1", size=11, color=ACCENT_BLUE, bold=True)
add_bullet(tf4r, "Shared ring buffer, batch I/O ops. ScyllaDB, TigerBeetle. Linux 5.1+", size=9, color=LIGHT_GRAY, level=1)
add_bullet(tf4r, "O_DIRECT \u2192 eliminates #2", size=11, color=ACCENT_BLUE, bold=True)
add_bullet(tf4r, "Bypass page cache, straight to user buffer. Needs 4K alignment.", size=9, color=LIGHT_GRAY, level=1)
add_bullet(tf4r, "posix_fadvise \u2192 helps #3", size=11, color=ACCENT_BLUE, bold=True)
add_bullet(tf4r, "SEQUENTIAL / WILLNEED hints. Nearly free.", size=9, color=LIGHT_GRAY, level=1)
add_bullet(tf4r, "Combined \u2192 eliminates all three", size=11, color=ACCENT_ORANGE, bold=True)
add_bullet(tf4r, "Theoretical ceiling for Linux I/O performance.", size=9, color=LIGHT_GRAY, level=1)

add_notes(s4, """SPEAKER NOTES - Slide 4: How Linux I/O Works & Where We Can Improve

User Buffer -- your application's memory where the KV cache data lives. A pinned PyTorch tensor. "Pinned" means the OS will never swap it to disk -- it stays in physical RAM.

Syscall -- when your program calls pread(), it can't talk to the storage device directly. It asks the kernel via a "context switch" -- CPU saves program state, switches to kernel mode, does the work, switches back. Like going through a security checkpoint every time. ~1-2 microseconds each. x1000 blocks = 1-2ms pure overhead.

Page Cache -- the kernel reads into its own memory first (page cache), then copies to your buffer. Like a library photocopying every book you request. For our workload (write once, read once) we're paying for a copy we'll never reuse.

Block Layer -- the kernel's I/O scheduler. Merges adjacent requests, reorders for optimal patterns. Problem: doesn't know our access pattern unless we tell it.

Device Driver / Storage -- NVMe drives have internal queues (up to 64K commands). The faster we fill this queue, the more the drive parallelizes. pread + threads: one request at a time per thread. io_uring: flood the queue with hundreds at once.

io_uring explained -- replaces per-operation syscalls with two shared memory ring buffers: submission queue (SQ) where you drop requests, completion queue (CQ) where kernel drops results. Fill SQ with hundreds of reads, notify kernel once. With SQPOLL mode, a kernel thread watches SQ continuously -- zero syscalls. Used by ScyllaDB, TigerBeetle. Linux 5.1+ (2019), mature.

O_DIRECT explained -- data goes storage -> your buffer directly, skipping page cache. Eliminates one full memcpy per block. Needs 4K-aligned buffers. Doesn't apply on tmpfs. Pure win for our read-once workload.

posix_fadvise explained -- one syscall hint after open(). SEQUENTIAL = aggressive readahead. WILLNEED = start prefetching now. Nearly free. Currently our code gives zero hints.

Why combined is the ceiling -- io_uring + O_DIRECT: no syscall overhead, no page cache copy, kernel optimally schedules. Only bottleneck left is the physical device.

Analogy: Current (pread): walk to security desk for each of 1000 boxes, each gets photocopied. io_uring + O_DIRECT: hand warehouse a list of all 1000 boxes, delivered directly to your truck. One trip, no copies.""")


# =============================================================================
# SLIDE 5 — Optimization Approaches
# =============================================================================
s5 = prs.slides.add_slide(BLANK)
make_title(s5, "Optimization Approaches")

# Section A
tb5ah = add_textbox(s5, L_MARGIN, Inches(0.9), CONTENT_W, Inches(0.3))
set_text(tb5ah.text_frame, "A. Optimize Current Approach (existing C++ backend)", size=13, color=ACCENT_GREEN, bold=True)

data5a = [
    ["Parameter", "What it does", "Effort"],
    ["O_NOATIME", "Skip access time metadata update on every read", "1 line"],
    ["posix_fadvise", "Tell kernel our access pattern (sequential / prefetch)", "2-3 lines"],
    ["I/O chunk size", "Tune pread/pwrite size within a block", "Small"],
    ["Prefetch depth", "fadvise(WILLNEED) on next N blocks while reading current", "Small"],
    ["Thread count range", "Currently hardcoded [16,32,64] \u2014 open to 8-256", "1 line"],
]
add_table(s5, 6, 3, L_MARGIN, Inches(1.25), CONTENT_W, Inches(1.5), data5a,
          col_widths=[Inches(1.8), Inches(5.6), Inches(1.4)])

tb5ag = add_textbox(s5, L_MARGIN, Inches(2.85), CONTENT_W, Inches(0.25))
set_text(tb5ag.text_frame, "Expected gain: ~5-15% with right combination. Minimal effort.", size=10, color=ACCENT_ORANGE, bold=True)

# Section B
tb5bh = add_textbox(s5, L_MARGIN, Inches(3.2), CONTENT_W, Inches(0.3))
set_text(tb5bh.text_frame, "B. New I/O Mechanisms", size=13, color=ACCENT_GREEN, bold=True)

data5b = [
    ["Approach", "What it does", "Expected read gain", "Effort"],
    ["io_uring", "Batch syscalls, kernel-side scheduling", "~10-20%", "New backend"],
    ["O_DIRECT", "Bypass page cache entirely", "~5-15%", "Flag + alignment"],
    ["io_uring + O_DIRECT", "All combined", "~15-30%", "New backend"],
]
add_table(s5, 4, 4, L_MARGIN, Inches(3.55), CONTENT_W, Inches(0.85), data5b,
          col_widths=[Inches(1.8), Inches(3.8), Inches(1.7), Inches(1.5)])

add_notes(s5, """SPEAKER NOTES - Slide 5: Optimization Approaches

Section A -- parameter explanations:

O_NOATIME -- every time Linux opens a file for reading, it updates the file's "last accessed" timestamp in filesystem metadata. A small write for every read -- useless for us. O_NOATIME skips this. One character in the open() call, zero downside.

posix_fadvise -- tell the kernel how we plan to read. SEQUENTIAL = aggressive readahead. WILLNEED = "start loading now." Currently our code gives zero hints -- the kernel is guessing.

I/O chunk size -- current code calls pread(fd, buffer, 50MB, offset). But different storage systems have different optimal transfer sizes. A GPFS filesystem with 4MB block size may prefer 4MB-aligned reads. Tunable chunk size lets us find the sweet spot.

Prefetch depth -- while reading block N, call fadvise(WILLNEED) on blocks N+1..N+K. Kernel starts loading them before we ask. K depends on storage speed vs processing speed.

Thread count range -- hardcoded [16, 32, 64]. Optimal depends on storage. NVMe might peak at 32, Spectrum Scale (striped across servers) might need 128-256. Slow device: 8 might beat 64.

Section B -- new mechanisms:

io_uring is what high-performance databases use (ScyllaDB, TigerBeetle). New backend but reuses existing benchmark framework. Also introduces its own tunable parameters (queue depth, SQPOLL, IOPOLL, registered buffers/files) -- roughly doubling our parameter space.

O_DIRECT caveat: needs 4K-aligned buffers, doesn't work on tmpfs. Must test on actual target storage.""")


# =============================================================================
# SLIDE 6 — Auto-Tuning & Roadmap
# =============================================================================
s6 = prs.slides.add_slide(BLANK)
make_title(s6, "Auto-Tuning & Roadmap")

# The problem
tb6p = add_textbox(s6, L_MARGIN, Inches(0.9), CONTENT_W, Inches(0.8))
tf6p = tb6p.text_frame
tf6p.word_wrap = True
set_text(tf6p, "The tuning problem:", size=13, color=ACCENT_BLUE, bold=True)
add_bullet(tf6p, "~15 tunable parameters across existing + new backends", size=11, color=WHITE)
add_bullet(tf6p, "Optimal values depend on hardware and interact non-trivially", size=11, color=WHITE)
add_bullet(tf6p, "Manual tuning doesn\u2019t scale across deployment targets", size=11, color=WHITE)

# Table
data6 = [
    ["Method", "What it tunes", "Best for"],
    ["Bayesian (Optuna)", "All parameters", "Per-machine optimal config"],
    ["Evolutionary (OpenEvolve)", "Parameters + code strategies", "Discovering new I/O patterns"],
]
add_table(s6, 3, 3, L_MARGIN, Inches(2.05), CONTENT_W, Inches(0.6), data6,
          col_widths=[Inches(2.4), Inches(3.0), Inches(3.4)])

# Roadmap
tb6r = add_textbox(s6, L_MARGIN, Inches(2.9), CONTENT_W, Inches(2.4))
tf6r = tb6r.text_frame
tf6r.word_wrap = True
set_text(tf6r, "Recommended Roadmap:", size=13, color=ACCENT_GREEN, bold=True)

steps = [
    ("1.", "Add tunable parameters + POSIX hints to existing C++ backend", "quick win"),
    ("2.", "Build io_uring + O_DIRECT backend", "main investment"),
    ("3.", "Bayesian auto-tuning per deployment target (Optuna)", "optimal config"),
    ("4.", "Evolutionary search for code-level strategies (OpenEvolve)", "exploration"),
]
for num, desc, tag in steps:
    p = tf6r.add_paragraph()
    p.space_before = Pt(3)
    p.space_after = Pt(1)
    r1 = p.add_run()
    r1.text = f"{num} "
    r1.font.size = Pt(11)
    r1.font.color.rgb = ACCENT_ORANGE
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = WHITE
    r3 = p.add_run()
    r3.text = f"  ({tag})"
    r3.font.size = Pt(9)
    r3.font.color.rgb = LIGHT_GRAY

add_paragraph(tf6r, "Build once, auto-tune everywhere.", size=12, color=ACCENT_ORANGE, bold=True, space_before=Pt(10))

add_notes(s6, """SPEAKER NOTES - Slide 6: Auto-Tuning & Roadmap

- Key insight: we don't need to understand every hardware detail. Build the parameter space, let the search find what works.

- Example of non-obvious interaction: O_DIRECT might help at 64 threads but hurt at 16. Or chunk size 1MB wins with SEQUENTIAL hint but 4MB wins with WILLNEED. You'd never find this manually.

- Bayesian optimization: finds best parameter values when each evaluation is expensive (like a benchmark run). After each trial, builds a statistical model ("surrogate") predicting which combinations look promising. Picks the most informative next trial -- balancing exploration vs exploitation. Much more efficient than grid/random search. Optuna is a popular Python library implementing this. You define parameter ranges + a function that returns throughput. Optuna handles the rest -- choosing what to try, tracking results, visualizing which parameters matter most.

- Evolutionary search (OpenEvolve) -- why code evolution, not just parameter tuning? Because some optimizations aren't expressible as numbers. Should each thread read one file or multiple files sequentially? Sort blocks by location or keep original order? Serialize file creation, pipeline in batches, or interleave with writes? These are branching code paths, not slider values. Bayesian can't explore them. Evolutionary search generates structurally different implementations, combines successful patterns, discovers strategies we wouldn't try manually. The current C++ backend is hand-optimized and good -- code evolution is how we break past incremental tuning ceilings.

- Roadmap: step 1 is low effort, enables all later work. Step 2 is core engineering. Step 3 makes it production-ready for any hardware. Step 4 is research/exploration.

- End with: "build once, auto-tune everywhere" -- same framework works whether target is local NVMe, Spectrum Scale, or future storage.""")


# =============================================================================
# Save
# =============================================================================
output_path = os.path.join(os.path.dirname(__file__), "io_optimization_presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
